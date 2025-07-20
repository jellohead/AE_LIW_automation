# slide_4.py

import logging
from pptx.util import Pt
from AE_LIW_automation.config import REPORTING_PERIOD, REPORTING_YEAR, CURRENT_MONTH_TEXT, CURRENT_YEAR
from AE_LIW_automation.helper_modules import update_paragraphs

logger = logging.getLogger(__name__)


# TODO: Rewrite to use update_paragraphs.py

def slide_4_updater(meta, df, df_labeled, prs) -> object:
    slide_index = 3
    print(
        f'\n================================\n======= Updating slide {slide_index + 1} =======\n================================\n')
    logger.info(f'Updating slide {slide_index + 1}')
    slide = prs.slides[slide_index]
    text_holder = None

    # extract existing text from textbox
    existing_text = ""
    for shape in slide.shapes:
        if shape.name == 'Rectangle 3' and shape.has_text_frame:
            text_frame = shape.text_frame
            existing_text = "\n".join([p.text for p in text_frame.paragraphs])
            break  # stop after finding the shape

    print("Extracted text:")
    print(existing_text)
    previous_qtr_percentage = existing_text.split('was ')[1].split(',')[0]
    print(f'{previous_qtr_percentage = }\n')
    previous_qtr_percentage_integer = int(previous_qtr_percentage.split('%')[0])


    q24_TopBox_result = (df['Q24']
                         .dropna()
                         .apply(lambda x: 'TopBox' if x in [8, 9, 10] else 'BottomBox')
                         .value_counts(normalize=True)
                         )['TopBox']

    q24_TopBox_result_integer = round(q24_TopBox_result * 100)

    if q24_TopBox_result_integer > previous_qtr_percentage_integer:
        comparison_description = 'up'
    elif q24_TopBox_result_integer < previous_qtr_percentage_integer:
        comparison_description = 'down'
    else:
        comparison_description = 'no change'

    q2_questions_list = [
        'Q2_1',
        'Q2_2',
        'Q2_3',
        'Q2_4',
        'Q2_5',
        'Q2_6',
        'Q2_7',
        'Q2_8',
        'Q2_9',
        'Q2_10',
        'Q2_11',
    ]

    q2_result_dict = {}
    # get top 3 responses to how customers found out about the weatherization program
    for question in q2_questions_list:
        dept_name = meta.column_names_to_labels[question].split('? ', 1)[1]
        result = df_labeled[question].dropna().value_counts().get('Checked', 0)
        q2_result_dict[dept_name] = result

    q2_result_dict_sorted = dict(sorted(q2_result_dict.items(), key=lambda x: x[1], reverse=True))
    print(q2_result_dict_sorted)
    q2_result_dict_sorted_top_3_keys = list(q2_result_dict_sorted.keys())[:3]

    paragraph_strings = [
        f'Overall satisfaction score with energy savings was {q24_TopBox_result:.0%} in {REPORTING_PERIOD} {REPORTING_YEAR}, {comparison_description} from {previous_qtr_percentage}%.',
        ' ',
        f'Contractor and customer service ratings remained relatively high for all attributes.',
        ' ',
        f'Customers appeared to be satisfied with the follow-up phone calls and indicated that the Austin Energy staff member/contractor did an overall good job on the work done at their homes.',
        ' ',
        f'For this quarter, {q2_result_dict_sorted_top_3_keys[0]}, {q2_result_dict_sorted_top_3_keys[1]}, and {q2_result_dict_sorted_top_3_keys[2]} were the top responses for how customers first learned about the weatherization program.',
        ' ',
        f'For this quarter, due to the small sample size, none of the changes can be deemed significant.'
        ]


    # text_holder = None
    # # get shape object by name
    # for shape in slide.shapes:
    #     if shape.name == 'Rectangle 3':
    #         shape.text_frame.clear()
    #         text_holder = shape.text_frame
    #
    # p = text_holder.paragraphs[0]
    #
    # for para_string in paragraph_strings:
    #     print(para_string)
    #     clean_text = para_string.replace('-', '')
    #     p = text_holder.add_paragraph()
    #     p.text = clean_text.strip()
    #     p.alignment = PP_ALIGN.LEFT
    #
    #     if para_string.startswith('-'):
    #         p.level = 1
    #     else:
    #         p.level = 0
    #
    #     run = p.runs[0]
    #     run.font.name = 'Tahoma'
    #     run.font.size = Pt(18 if p.level == 0 else 16)

    # revise to use update_paragraphs.py
    for shape in slide.shapes:
        if shape.name == 'Rectangle 3':
            shape.text_frame.clear()
            text_holder = shape.text_frame
            break

    update_paragraphs(paragraph_strings,
                      text_holder,
                      l0_font_size=Pt(18)
                      )

