# slide_8.py
# update the data table on slide 8

# TODO: Refactor using table helper functions
# TODO: Break style_cell function into a module
# TODO: Break style_cell_old_text function into a module
# TODO: Fix styling to center text vertically in cells
import logging
from pandas import DataFrame
import pandas as pd
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from AE_LIW_automation.helper_modules import get_table_object
from AE_LIW_automation.config import REPORTING_PERIOD, REPORTING_YEAR, CURRENT_MONTH_TEXT, CURRENT_YEAR


logger = logging.getLogger(__name__)

# Function to apply styling to a table cell
def style_cell(cell, text, font_size=12, bold=False, color=RGBColor(0, 0, 0), bg_color=None, align=PP_ALIGN.CENTER):
    """Applies font size, boldness, color, and background to a cell."""
    cell.text = text
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.alignment = align
    if paragraph.runs:
        run = paragraph.runs[0]
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color

# Function to apply styling to a table cell with existing text
def style_cell_old_text(cell, font_size=12, bold=False, color=RGBColor(0, 0, 0), bg_color=None, align=PP_ALIGN.CENTER):
    """Applies font size, boldness, color, and background to a cell."""
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.alignment = align
    if paragraph.runs:
        run = paragraph.runs[0]
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color

def slide_8_updater(df, prs):
    slide_index = 7
    print(
        f'\n================================\n======= Updating slide {slide_index + 1} =======\n================================\n')
    logger.info(f'Updating slide {slide_index + 1}')

    slide = prs.slides[slide_index]

    table = get_table_object(slide)
    if not table:
        print(f'No table found on {slide_index + 1}')
        return

    # extract table dimensions
    num_rows = len(table.rows)
    num_cols = len(table.columns)

    # Get existing data from old table
    table_data = [[cell.text.strip() for cell in row.cells] for row in table.rows]
    table_df: DataFrame = pd.DataFrame(table_data[1:], columns=table_data[0])
    table_df.set_index(table_df.columns[0], inplace=True)

    # drop oldest quarter data
    table_df_current = table_df.drop(columns=[table_df.columns[0]])
    print(f'{table_df_current = }')

    # get current quarter data from dataset
    q19_result = df['Q19'].replace('', float('nan')).dropna().value_counts()
    # append Base value to the series
    q19_result.at['Base:'] = len(q19_result)
    q19_df =pd.DataFrame({f'{REPORTING_PERIOD} {REPORTING_YEAR}': q19_result}).fillna(0)

    # combine old and new data, convert new data to integers
    q19_df_combined = pd.concat([table_df_current, q19_df], axis=1).fillna(0).astype(int)

    # sort the combined df and put Base at the bottom
    base_row = q19_df_combined[q19_df_combined.index == 'Base:']
    print(f'{type(base_row) = }')
    not_base_rows = q19_df_combined[q19_df_combined.index != 'Base:'].sort_values(by=f'{REPORTING_PERIOD} {REPORTING_YEAR}', ascending=False)
    q19_df_combined = pd.concat([not_base_rows, base_row], axis=0).fillna(0).astype(int)    # q19_df_other_rows = q19_df_combined.loc[q19_df_combined.index != 'Base:'].sort_values(f'{REPORTING_PERIOD} {REPORTING_YEAR}', ascending=False)

    # Step 1: Remove existing table (if any)
    shapes = slide.shapes
    for shape in shapes:
        if shape.has_table:  # Check if shape is a table
            sp = shape._element  # Get the XML element of the shape
            slide.shapes._spTree.remove(sp)  # Remove the shape

    # Step 3: Add a new table to the slide
    rows, cols = q19_df_combined.shape

    # Define styling properties
    header_bg_color = RGBColor(90,	128,184)  # Dark Blue
    header_text_color = RGBColor(255, 255, 255)  # White
    row_bg_color = RGBColor(224, 235, 255)  # Light Blue (Alternating Rows)
    last_row_bg_color = RGBColor(90,	128,	184)  # Dark Blue for Last Row
    data_text_color = RGBColor(0, 0, 0) # Black text
    data_bg_color = RGBColor(224,	229,	240) # Light blue for data rows



    # Add one more column for the index
    table_shape = slide.shapes.add_table(rows + 1, cols + 1, Inches(.5), Inches(1.7), Inches(6.5), Inches(5)).table

    # Step 4: Insert column headers (including index column)
    # add styling to the header row
    style_cell(table_shape.cell(0, 0), text='',font_size=12, bold=True, color=header_text_color, bg_color=header_bg_color, align=PP_ALIGN.CENTER)
    for col_idx, col_name in enumerate(q19_df_combined.columns):
        style_cell(table_shape.cell(0, col_idx + 1), col_name, font_size=14, bold=True, color=header_text_color, bg_color=header_bg_color, align=PP_ALIGN.CENTER)

    # Step 5: Insert data rows (including index values)
    for row_idx, (index_value, row) in enumerate(q19_df_combined.iterrows()):
        # table_shape.cell(row_idx + 1, 0).text = str(index_value)  # First column for index
        style_cell(table_shape.cell(row_idx + 1, 0),
                   str(index_value),
                   font_size=13,
                   bold=False,
                   color=data_text_color,
                   bg_color=data_bg_color,
                   align=PP_ALIGN.LEFT)
        for col_idx, value in enumerate(row):
            # table_shape.cell(row_idx + 1, col_idx + 1).text = str(value)  # Shifted by +1
            style_cell(table_shape.cell(row_idx + 1, col_idx + 1),
                       str(value),
                       font_size=13,
                       bold=False,
                       color=data_text_color,
                       bg_color=data_bg_color,
                       align=PP_ALIGN.CENTER)

    # TODO pull data from the correct location for last row
    # add styling to the last row of the table
    style_cell(table_shape.cell(len((q19_df_combined)), 0), text=base_row.index[0], font_size=12, bold=True, color=header_text_color,
               bg_color=header_bg_color, align=PP_ALIGN.LEFT)


    for col_number, value in enumerate(q19_df_combined.loc['Base:']):
        style_cell_old_text(table_shape.cell(len(q19_df_combined), col_number + 1),
                                             font_size=13,
                                             bold=False,
                                             color=header_text_color,
                                             bg_color=header_bg_color,
                                             align=PP_ALIGN.CENTER),
        print(f'{value = }')

    logger.info(f'Update of slide {slide_index + 1} complete.\nManually remove rows where all values are zero.\nManually adjust position and size of the table.')
