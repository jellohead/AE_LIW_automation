# slide_3.py

from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

from AE_LIW_automation.config import REPORTING_PERIOD, REPORTING_YEAR, CURRENT_MONTH_TEXT, CURRENT_YEAR


def slide_3_updater(df, prs) -> object:
    print('updating slide 3')
    slide_index = 2
    slide = prs.slides[slide_index]

    q11_topbox_result = (df['Q11']
                         .dropna()
                         .apply(lambda x: 'TopBox'if x in [8, 9, 10] else 'BottomBox')
                         .value_counts(normalize=True)
                         )['TopBox']

    q31_topbox_result = (df['Q31']
                         .dropna()
                         .apply(lambda x: 'TopBox'if x in [8, 9, 10] else 'BottomBox')
                         .value_counts(normalize=True)
                         )['TopBox']

    q15_result = (df['Q15']
                  .dropna()
                  .value_counts(normalize=True)
                  )[1]

    q22_topbox_result = (df['Q22']
                         .dropna()
                         .apply(lambda x: 'TopBox'if x in [8, 9, 10] else 'BottomBox')
                         .value_counts(normalize=True)
                         )['TopBox']

    paragraph_strings = [f'Customers’ overall satisfaction with Austin Energy’s Weatherization Program remained favorable with a score of {q11_topbox_result:.0%} for {REPORTING_PERIOD} {REPORTING_YEAR}.',
                         f'Overall satisfaction level with Austin Energy improved to {q31_topbox_result:.0%} in {REPORTING_PERIOD} {REPORTING_YEAR} from 89% in Q2 2024.',
                         f'Customers indicated a need for the program and home weatherization assistance.',
                         f'{q15_result:.0%} of customers would recommend this program to a friend or family member.',
                         f'{q22_topbox_result:.0%} of customers appeared to have a high level of understanding when it comes to their utility bill and energy savings.'
                         ]

    text_holder = None
    # get shape object by name
    for shape in slide.shapes:
        if shape.name == 'Rectangle 3':
            shape.text_frame.clear()
            text_holder = shape.text_frame

    p = text_holder.paragraphs[0]

    for para_string in paragraph_strings:
        print(para_string)
        clean_text = para_string.replace('-', '')
        p = text_holder.add_paragraph()
        p.text = clean_text.strip()
        p.alignment = PP_ALIGN.LEFT

        if para_string.startswith('-'):
            p.level = 1
        else:
            p.level = 0


        run = p.runs[0]
        run.font.name = 'Tahoma'
        run.font.size = Pt(18 if p.level == 0 else 16)
