# slide_33.py

import logging
from pandas import DataFrame, Series
import pandas as pd
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from AE_LIW_automation.helper_modules import get_table_object, style_table_cell, combine_multiple_questions
from AE_LIW_automation.config import REPORTING_PERIOD, REPORTING_YEAR


logger = logging.getLogger(__name__)


# TODO: Refactor using table helper functions
# TODO: v_alignment is not modifying the vertical alignment of the cell text
# TODO: write helper module to combine results across multiple questions

def slide_33_updater(meta, df, df_labeled, prs):
    slide_index = 32
    print(
        f'\n================================\n======= Updating slide {slide_index + 1} =======\n================================\n')
    logger.info(f'Updating slide {slide_index + 1}')

    slide = prs.slides[slide_index]

    question_list = ['Q13_1', 'Q13_2', 'Q13_3', 'Q13_4',]
    label_sub_dict = {'All other' : 'Other',
                      'Do not remember, do not know': "Don't know",
                      'Nothing': 'Nothing/no changes',
                      }
    last_rows = ["Don't know", 'Other', 'Base:']

    table = get_table_object(slide)
    if not table:
        print(f'No table found on {slide_index + 1}')
        logger.info(f'No table found on {slide_index + 1}')
        return

    # Get existing data from old table
    table_data = [[cell.text.strip() for cell in row.cells] for row in table.rows]
    table_df: DataFrame = pd.DataFrame(table_data[1:], columns=table_data[0])
    table_df.set_index(table_df.columns[0], inplace=True)

    # drop oldest quarter data
    table_df_existing = table_df.drop(columns=[table_df.columns[0]])
    print(f'{table_df_existing = }')


    # get current quarter data from dataset
    current_quarter_result_series = combine_multiple_questions(df_labeled, question_list, label_sub_dict)
    current_quarter_result_df = pd.DataFrame({f'{REPORTING_PERIOD} {REPORTING_YEAR}': current_quarter_result_series}).fillna(0)

    # combine old and new data, convert new data to integers
    current_quarter_result_df_combined = pd.concat([table_df_existing, current_quarter_result_df], axis=1).fillna(
        0).astype(int)

    # remove rows where all values are 0
    current_quarter_result_df_combined = current_quarter_result_df_combined[
        ~(current_quarter_result_df_combined == 0)
        .all(axis=1)
    ]

    # sort the combined df and put last_rows at the bottom in correct order
    rows_to_move = {label: current_quarter_result_df_combined[current_quarter_result_df_combined.index == label]
                    for label in last_rows
                    }

    remaining_rows = current_quarter_result_df_combined[~current_quarter_result_df_combined.index.isin(last_rows)]
    remaining_rows_sorted = (remaining_rows.sort_values(by=f'{REPORTING_PERIOD} {REPORTING_YEAR}',ascending=False))
    ordered_rows = [rows_to_move[label] for label in last_rows if not rows_to_move[label].empty]

    current_quarter_result_df_combined = pd.concat(
        [remaining_rows_sorted] + ordered_rows,
        axis=0
    ).fillna(0).astype(int)

    base_row = current_quarter_result_df_combined[current_quarter_result_df_combined.index == 'Base:']

    # Step 1: Remove existing table (if any)
    shapes = slide.shapes
    for shape in shapes:
        if shape.has_table:  # Check if shape is a table
            sp = shape._element  # Get the XML element of the shape
            slide.shapes._spTree.remove(sp)  # Remove the shape

    # Step 3: Add a new table to the slide
    rows, cols = current_quarter_result_df_combined.shape

    # Define styling properties
    header_bg_color = RGBColor(90, 128, 184)  # Dark Blue
    header_text_color = RGBColor(255, 255, 255)  # White
    row_bg_color = RGBColor(224, 235, 255)  # Light Blue (Alternating Rows)
    last_row_bg_color = RGBColor(90, 128, 184)  # Dark Blue for Last Row
    data_text_color = RGBColor(0, 0, 0)  # Black text
    data_bg_color = RGBColor(224, 229, 240)  # Light blue for data rows

    # Add one more column for the index
    table_shape = slide.shapes.add_table(rows + 1, cols + 1, Inches(.5), Inches(1.7), Inches(6.5), Inches(5)).table

    # Step 4: Insert column headers (including index column)
    # add styling to the header row
    style_table_cell(table_shape.cell(0, 0),
                     text='',
                     font_size=12,
                     bold=True,
                     color=header_text_color,
                     bg_color=header_bg_color,
                     h_alignment=PP_ALIGN.CENTER,
                     v_alignment=MSO_ANCHOR.MIDDLE,
                     )
    for col_idx, col_name in enumerate(current_quarter_result_df_combined.columns):
        style_table_cell(table_shape.cell(0, col_idx + 1),
                         col_name,
                         font_size=14,
                         bold=True,
                         color=header_text_color,
                         bg_color=header_bg_color,
                         h_alignment=PP_ALIGN.CENTER,
                         v_alignment=MSO_ANCHOR.MIDDLE,
                         )

    # Step 5: Insert data rows (including index values)
    for row_idx, (index_value, row) in enumerate(current_quarter_result_df_combined.iterrows()):
        style_table_cell(table_shape.cell(row_idx + 1, 0),
                         str(index_value),
                         font_size=13,
                         bold=False,
                         color=data_text_color,
                         bg_color=data_bg_color,
                         h_alignment=PP_ALIGN.LEFT,
                         v_alignment=MSO_ANCHOR.MIDDLE,
                         )
        for col_idx, value in enumerate(row):
            style_table_cell(table_shape.cell(row_idx + 1, col_idx + 1),
                             str(value),
                             font_size=13,
                             bold=False,
                             color=data_text_color,
                             bg_color=data_bg_color,
                             h_alignment=PP_ALIGN.CENTER,
                             v_alignment=MSO_ANCHOR.MIDDLE,
                             )

    # add styling to the last row of the table
    style_table_cell(table_shape.cell(len((current_quarter_result_df_combined)), 0),
                     text=base_row.index[0],
                     font_size=12,
                     bold=True,
                     color=header_text_color,
                     bg_color=header_bg_color,
                     h_alignment=PP_ALIGN.LEFT,
                     v_alignment=MSO_ANCHOR.MIDDLE,
                     )

    for col_number, value in enumerate(current_quarter_result_df_combined.loc['Base:']):
        style_table_cell(table_shape.cell(len(current_quarter_result_df_combined), col_number + 1),
                         font_size=13,
                         bold=False,
                         color=header_text_color,
                         bg_color=header_bg_color,
                         h_alignment=PP_ALIGN.CENTER,
                         v_alignment=MSO_ANCHOR.MIDDLE,
                         )

    logger.info(
        f'Update of slide {slide_index + 1} complete.\nManually adjust position and size of the table.\nVerify Base value for current quarter is accurate.')
