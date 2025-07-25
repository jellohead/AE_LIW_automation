# slide_52.py

import logging
from pandas import DataFrame, Series
import pandas as pd
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from AE_LIW_automation.helper_modules import get_table_object, style_table_cell, combine_multiple_questions
from AE_LIW_automation.config import REPORTING_PERIOD, REPORTING_YEAR

logger = logging.getLogger(__name__)


# TODO: recode no/none/nothing responses to group them


def slide_52_updater(meta, df, df_labeled, prs):
    slide_index = 51
    print(
        f'\n================================\n======= Updating slide {slide_index + 1} =======\n================================\n')
    logger.info(f'Updating slide {slide_index + 1}')

    slide = prs.slides[slide_index]

    question = 'D12'
    new_quarter_label = f'{REPORTING_PERIOD} {REPORTING_YEAR}'
    label_sub_dict = {}
    # last_rows = ['All other', "Don’t know", "Don't know", 'Base:']
    last_rows = ['All other', "Don’t know", 'Base:']

    table = get_table_object(slide)
    if not table:
        print(f'No table found on {slide_index + 1}')
        logger.info(f'No table found on {slide_index + 1}')
        return

    # Get existing data from old table
    table_data = [[cell.text.strip() for cell in row.cells] for row in table.rows]
    table_df: DataFrame = pd.DataFrame(table_data[1:], columns=table_data[0])
    table_df.set_index(table_df.columns[0], inplace=True)

    # drop oldest quarter data
    table_df_existing = table_df.drop(columns=[table_df.columns[0]])

    # get the existing base row and then drop it from the working dataframe
    base_row_df = table_df_existing.loc[['Base:']]
    base_row_df[new_quarter_label] = len(df)

    # get current quarter data from dataset
    current_quarter_result_series = df_labeled[question].value_counts(normalize=True).map("{:.0%}".format)
    current_quarter_result_df = pd.DataFrame({new_quarter_label: current_quarter_result_series})
    current_quarter_result_df.loc['Base:'] = len(df)
    current_quarter_result_df_combined = pd.concat([table_df_existing, current_quarter_result_df], axis=1).fillna('0%')

    # remove rows where all values are 0
    current_quarter_result_df_combined = current_quarter_result_df_combined[
        ~(current_quarter_result_df_combined == '0%')
        .all(axis=1)
    ]

    rows_to_move_df = current_quarter_result_df_combined[current_quarter_result_df_combined.index.isin(last_rows)]
    remaining_rows_df = current_quarter_result_df_combined[~current_quarter_result_df_combined.index.isin(last_rows)].copy()

    # sort dataframe by the current quarter column
    last_col = remaining_rows_df.columns[-1]
    remaining_rows_df[last_col] = remaining_rows_df[last_col].str.rstrip('%').astype(float)
    remaining_rows_sorted = remaining_rows_df.sort_values(by=last_col, ascending=False).copy()

    # reapply % format to last column, change dtype of last column to string
    remaining_rows_sorted = remaining_rows_sorted.astype({last_col: 'object'})
    remaining_rows_sorted.loc[:, last_col] = remaining_rows_sorted[last_col].map(lambda x: f'{x:.0f}%')


    final_result_df_combined = pd.concat(
        [remaining_rows_sorted, rows_to_move_df]
    )

    base_row = final_result_df_combined[final_result_df_combined.index == 'Base:']

    # Step 1: Remove existing table (if any)
    shapes = slide.shapes
    for shape in shapes:
        if shape.has_table:  # Check if shape is a table
            sp = shape._element  # Get the XML element of the shape
            slide.shapes._spTree.remove(sp)  # Remove the shape

    # Step 3: Add a new table to the slide
    rows, cols = current_quarter_result_df_combined.shape

    # Define styling properties
    header_bg_color = RGBColor(90, 128, 184)  # Dark Blue
    header_text_color = RGBColor(255, 255, 255)  # White
    row_bg_color = RGBColor(224, 235, 255)  # Light Blue (Alternating Rows)
    last_row_bg_color = RGBColor(90, 128, 184)  # Dark Blue for Last Row
    data_text_color = RGBColor(0, 0, 0)  # Black text
    data_bg_color = RGBColor(224, 229, 240)  # Light blue for data rows

    # Add one more column for the index
    table_shape = slide.shapes.add_table(rows + 1, cols + 1, Inches(.5), Inches(1.7), Inches(6.5), Inches(5)).table

    # Step 4: Insert column headers (including index column)
    # add styling to the header row
    style_table_cell(table_shape.cell(0, 0),
                     text='',
                     font_size=12,
                     bold=True,
                     color=header_text_color,
                     bg_color=header_bg_color,
                     h_alignment=PP_ALIGN.CENTER,
                     v_alignment=MSO_ANCHOR.MIDDLE,
                     )
    for col_idx, col_name in enumerate(final_result_df_combined.columns):
        style_table_cell(table_shape.cell(0, col_idx + 1),
                         col_name,
                         font_size=12,
                         bold=True,
                         color=header_text_color,
                         bg_color=header_bg_color,
                         h_alignment=PP_ALIGN.CENTER,
                         v_alignment=MSO_ANCHOR.MIDDLE,
                         )

    # Step 5: Insert data rows (including index values)
    for row_idx, (index_value, row) in enumerate(final_result_df_combined.iterrows()):
        style_table_cell(table_shape.cell(row_idx + 1, 0),
                         str(index_value),
                         font_size=12,
                         bold=False,
                         color=data_text_color,
                         bg_color=data_bg_color,
                         h_alignment=PP_ALIGN.LEFT,
                         v_alignment=MSO_ANCHOR.MIDDLE,
                         )
        for col_idx, value in enumerate(row):
            style_table_cell(table_shape.cell(row_idx + 1, col_idx + 1),
                             str(value),
                             font_size=12,
                             bold=False,
                             color=data_text_color,
                             bg_color=data_bg_color,
                             h_alignment=PP_ALIGN.CENTER,
                             v_alignment=MSO_ANCHOR.MIDDLE,
                             )

    # add styling to the last row of the table
    style_table_cell(table_shape.cell(len(final_result_df_combined), 0),
                     text=base_row.index[0],
                     font_size=12,
                     bold=True,
                     color=header_text_color,
                     bg_color=header_bg_color,
                     h_alignment=PP_ALIGN.LEFT,
                     v_alignment=MSO_ANCHOR.MIDDLE,
                     )

    for col_number, value in enumerate(final_result_df_combined.loc['Base:']):
        style_table_cell(table_shape.cell(len(final_result_df_combined), col_number + 1),
                         font_size=12,
                         bold=True,
                         color=header_text_color,
                         bg_color=header_bg_color,
                         h_alignment=PP_ALIGN.CENTER,
                         v_alignment=MSO_ANCHOR.MIDDLE,
                         )

    logger.info(
        f'Update of slide {slide_index + 1} complete.\nCopy/paste these results to populate slides 51 and 52.\nManually adjust position and size of the table.')
