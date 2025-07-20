# slide_56.py

import typing
import logging
from pandas import DataFrame, Series
import pandas as pd
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from AE_LIW_automation.helper_modules import get_table_object, style_table_cell, combine_multiple_questions, \
    get_table_shape_by_name
from AE_LIW_automation.config import REPORTING_PERIOD, REPORTING_YEAR

logger = logging.getLogger(__name__)


# TODO: recode no/none/nothing responses to group them
# TODO: add logic to handle if table requires sorting by values


def slide_56_updater(meta, df, df_labeled, prs):
    slide_index = 55
    print(
        f'\n================================\n======= Updating slide {slide_index + 1} =======\n================================\n')
    logger.info(f'Updating slide {slide_index + 1}')

    slide = prs.slides[slide_index]

    question_dict = {
                    'Table 1':'D6',
                     'Table 5': 'D7'}

    upper_table = list(question_dict.keys())[0]

    new_quarter_label = f'{REPORTING_PERIOD} {REPORTING_YEAR}'
    label_sub_dict = {'Prefer not to respond': 'Refused',
                      'DK/unsure': 'Do not know/unsure'}
    last_rows: list[str] = ['Do not know/unsure','Other', 'Refused', 'Base:']
    # table_names: list[str] = ['Table 1', 'Table 2']

    for table_name, question in question_dict.items():
        table_shape = get_table_shape_by_name(slide, table_name)
        if not table_shape:
            print(f'No table found on {slide_index + 1}')
            logger.info(f'No table found on {slide_index + 1}')
            return

        old_table_name = table_shape.name
        logger.info(f'Updating "{old_table_name}"')
        table = table_shape.table

        # get a reference to the table shape for the old table
        table_shape_element = table_shape._element

        # Get existing data from old table
        table_data = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        table_df: DataFrame = pd.DataFrame(table_data[1:], columns=table_data[0])
        table_df.set_index(table_df.columns[0], inplace=True)

        # drop oldest quarter data
        table_df_existing = table_df.drop(columns=[table_df.columns[0]])

        index_name = table_df_existing.index.name

        # get the existing base row and then drop it from the working dataframe
        base_row_df = table_df_existing.loc[['Base:']]
        base_row_df[new_quarter_label] = len(df)

        # get current quarter data from dataset
        current_quarter_result_series = df_labeled[question].value_counts(normalize=True).map("{:.0%}".format)
        current_quarter_result_df = pd.DataFrame({new_quarter_label: current_quarter_result_series})
        current_quarter_result_df.loc['Base:'] = len(df)
        current_quarter_result_df_combined = pd.concat([table_df_existing, current_quarter_result_df], axis=1).fillna('0%')
        current_quarter_result_df_combined.rename(index=label_sub_dict, inplace=True)

        # remove rows where all values are 0
        current_quarter_result_df_combined = current_quarter_result_df_combined[
            ~(current_quarter_result_df_combined == '0%')
            .all(axis=1)
        ]

        rows_to_move_df = current_quarter_result_df_combined[current_quarter_result_df_combined.index.isin(last_rows)]
        remaining_rows_df = current_quarter_result_df_combined[~current_quarter_result_df_combined.index.isin(last_rows)].copy()

        final_result_df_combined = pd.concat(
            [remaining_rows_df, rows_to_move_df]
        )

        final_result_df_combined.rename_axis(index_name, inplace=True)

        base_row = final_result_df_combined[final_result_df_combined.index == 'Base:']

        # Step 1: Remove existing table (if any)
        slide.shapes._spTree.remove(table_shape_element)

        rows, cols = current_quarter_result_df_combined.shape

        # Define styling properties
        header_bg_color = RGBColor(90, 128, 184)  # Dark Blue
        header_text_color = RGBColor(255, 255, 255)  # White
        row_bg_color = RGBColor(224, 235, 255)  # Light Blue (Alternating Rows)
        last_row_bg_color = RGBColor(90, 128, 184)  # Dark Blue for Last Row
        data_text_color = RGBColor(0, 0, 0)  # Black text
        data_bg_color = RGBColor(224, 229, 240)  # Light blue for data rows

        # provide a horizontal offset so tables do not overlap in case of two tables being updated
        top_offset = Inches(1.7) if table_name == upper_table else Inches(4.5)

        # Step 3: Add a new table to the slide
        # table_shape = slide.shapes.add_table(rows + 1, cols + 1, Inches(.5), top_offset, Inches(6.5), Inches(2.5)).table
        table_shape_object = slide.shapes.add_table(rows + 1, cols + 1, Inches(.5), top_offset, Inches(6.5), Inches(2.5))
        # rename new table to match old table name
        table_shape_object.name = old_table_name

        # get reference to table contained in the table shape object
        table_shape = table_shape_object.table

        # Step 4: Insert column headers (including index column)
        # add styling to the header row
        style_table_cell(table_shape.cell(0, 0),
                         text=index_name,
                         font_size=12,
                         bold=True,
                         color=header_text_color,
                         bg_color=header_bg_color,
                         h_alignment=PP_ALIGN.CENTER,
                         v_alignment=MSO_ANCHOR.MIDDLE,
                         )
        for col_idx, col_name in enumerate(final_result_df_combined.columns):
            style_table_cell(table_shape.cell(0, col_idx + 1),
                             col_name,
                             font_size=12,
                             bold=True,
                             color=header_text_color,
                             bg_color=header_bg_color,
                             h_alignment=PP_ALIGN.CENTER,
                             v_alignment=MSO_ANCHOR.MIDDLE,
                             )

        # Step 5: Insert data rows (including index values)
        for row_idx, (index_value, row) in enumerate(final_result_df_combined.iterrows()):
            style_table_cell(table_shape.cell(row_idx + 1, 0),
                             str(index_value),
                             font_size=12,
                             bold=False,
                             color=data_text_color,
                             bg_color=data_bg_color,
                             h_alignment=PP_ALIGN.LEFT,
                             v_alignment=MSO_ANCHOR.MIDDLE,
                             )
            for col_idx, value in enumerate(row):
                style_table_cell(table_shape.cell(row_idx + 1, col_idx + 1),
                                 str(value),
                                 font_size=12,
                                 bold=False,
                                 color=data_text_color,
                                 bg_color=data_bg_color,
                                 h_alignment=PP_ALIGN.CENTER,
                                 v_alignment=MSO_ANCHOR.MIDDLE,
                                 )

        # add styling to the last row of the table
        style_table_cell(table_shape.cell(len((final_result_df_combined)), 0),
                         text=base_row.index[0],
                         font_size=12,
                         bold=True,
                         color=header_text_color,
                         bg_color=header_bg_color,
                         h_alignment=PP_ALIGN.LEFT,
                         v_alignment=MSO_ANCHOR.MIDDLE,
                         )

        for col_number, value in enumerate(final_result_df_combined.loc['Base:']):
            style_table_cell(table_shape.cell(len(final_result_df_combined), col_number + 1),
                             font_size=12,
                             bold=True,
                             color=header_text_color,
                             bg_color=header_bg_color,
                             h_alignment=PP_ALIGN.CENTER,
                             v_alignment=MSO_ANCHOR.MIDDLE,
                             )

    logger.info(
        f'Update of slide {slide_index + 1} complete.\nManually adjust position and size of the table.')
