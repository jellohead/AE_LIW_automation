# slide_1.py
# This file contains the functions for updating the slide 1 of the Powerpoint file.


from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
# from config import REPORTING_PERIOD, REPORTING_YEAR, CURRENT_MONTH_TEXT, CURRENT_YEAR

from config import REPORTING_PERIOD, REPORTING_YEAR, CURRENT_MONTH_TEXT, CURRENT_YEAR
from helper_modules import get_chart_object_by_name

# TODO Text color is wrong

def slide_1_updater(df, prs) -> object:
    # global text_holder
    print('updating slide 1')
    slide_index = 0
    slide = prs.slides[slide_index]

    paragraph_strings = [f'Low Income Weatherization Survey\nFY{REPORTING_YEAR} {REPORTING_PERIOD}\n\n',
                         f'-{CURRENT_MONTH_TEXT} {CURRENT_YEAR}'
                         ]

    text_holder = None
    # get shape object by name
    for shape in slide.shapes:
        if shape.name == 'Rectangle 2':
            shape.text_frame.clear()
            text_holder = shape.text_frame

    p = text_holder.paragraphs[0]

    for para_string in paragraph_strings:
        print(para_string)
        run = p.add_run()
        if para_string.startswith('-'):
            run.text = para_string.replace ('-', '')
            run.font.name = 'Arial'
            run.font.size = Pt(16)
            run.alignment = PP_ALIGN.CENTER
        else:
            run.text = para_string
            run.font.name = 'Arial'
            run.font.size = Pt(28)
            run.alignment = PP_ALIGN.CENTER