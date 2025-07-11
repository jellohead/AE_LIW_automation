from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Pt

# TODO: v_alignment is not working

def style_table_cell(cell,
                     text=None,
                     font_name='Tahoma',
                     font_size=12,
                     bold=False,
                     italic=False,
                     color=RGBColor(0, 0, 0), # black is default
                     bg_color=None,
                     h_alignment=PP_ALIGN.CENTER,
                     v_alignment=MSO_ANCHOR.MIDDLE,
                     ):
    """
    Applies font size, boldness, color, alignment, and background to a table cell.

    Parameters:
        cell      : PowerPoint table cell object
        text      : Optional; if provided, sets the cell text
        font_name : Optional; if provided, sets the cell font name
        font_size : Font size in points
        bold      : Boolean, whether font is bold
        italic    : Boolean, whether font is italic
        color     : RGBColor instance for font color, default black
        bg_color  : Optional RGBColor for background fill
        h_alignment     : Paragraph horizontal alignment (e.g., PP_ALIGN.CENTER)
        v_alignment    : Paragraph vertical alignment (e.g., MSO_ANCHOR.MIDDLE)
    """



    if text is not None:
        cell.text = text

    # cell.text_frame.vertical_anchor = v_alignment (referencing text_frame for table cells does not work)
    cell.vertical_anchor = v_alignment # do this instead

    paragraph = cell.text_frame.paragraphs[0]
    paragraph.alignment = h_alignment

    if paragraph.runs:
        run = paragraph.runs[0]
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color

    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color

